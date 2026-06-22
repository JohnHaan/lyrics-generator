"""
raw/ 에 추가된 두 번째 가사 묶음(9곡)을 표준 형식으로 변환해 songs/song-11~19로
등록한다. 이 묶음은 첫 번째 묶음과 달리:

- 제목 슬라이드가 없음 -> 파일명(확장자 제외)으로 제목 슬라이드를 새로 만든다.
- 폰트가 테마 기본값("맑은 고딕" 60pt, schemeClr 기반 색)이라 표준
  (휴먼모음T, 72pt, 흰색/검정외곽선)과 다름 -> 기존 슬라이드 XML을 고치는 대신
  텍스트만 추출해서 표준 스타일로 슬라이드를 새로 만든다 (ooxml_utils의
  lyric_slide_from_lines 사용).
- 가사 내용 자체는 수정하지 않는다("X2", "X∞" 같은 반복 표시도 원문 그대로 유지).
  단, 사용자가 명시적으로 승인한 2건만 예외로 제거한다:
    1) "사랑한다 말하시네" 슬라이드6의 가사 아닌 안내문 "(다음 이어가기)" 한 줄 제거
    2) "태산을 넘어 험곡에 가도" 슬라이드3("간주"만 있는 진행 표시 슬라이드) 전체 제거
- key는 라벨을 새로 짓지 않고 "slideN"처럼 등장 순서 기반으로 부여하고,
  완전히 동일한 텍스트의 슬라이드만 같은 key로 합친다(라벨 추측/병합 없음).
"""
import json
from pathlib import Path

from ooxml_utils import build_song_pptx, lyric_slide_from_lines, title_slide
from pptx import Presentation

ROOT = Path(__file__).resolve().parent.parent
RAW_DIR = ROOT / "raw"
TEMPLATE = ROOT / "template" / "blank.pptx"
SONGS_DIR = ROOT / "songs"


def lines_of(text: str) -> list[str]:
    return text.replace("\x0b", "\n").split("\n")


# 각 곡: 파일명, slug, (선택) 슬라이드 텍스트 전처리/제거 규칙.
# "drop"은 0-based 원본 슬라이드 인덱스를 제거, "strip_last_line"은 마지막 줄만 제거.
BATCH2 = [
    {"file": "나 무엇과도 주님을 바꾸지 않으리.pptx", "slug": "song-11"},
    {"file": "내 맘의 눈을 여소서.pptx", "slug": "song-12"},
    {"file": "내 이름 아시죠.pptx", "slug": "song-13"},
    {"file": "내가 어둠 속에서.pptx", "slug": "song-14"},
    {"file": "사랑의 노래되리.pptx", "slug": "song-15"},
    {
        "file": "사랑한다 말하시네.pptx",
        "slug": "song-16",
        "strip_last_line": {5},  # 0-based: 슬라이드6에서 "(다음 이어가기)" 제거
    },
    {"file": "주께서 주신 동산에.pptx", "slug": "song-17"},
    {
        "file": "태산을 넘어 험곡에 가도.pptx",
        "slug": "song-18",
        "drop": {2},  # 0-based: 슬라이드3("간주") 전체 제거
    },
    {"file": "하나님은 너를 지키시는 자.pptx", "slug": "song-19"},
]


def main() -> None:
    for song in BATCH2:
        title = Path(song["file"]).stem
        prs = Presentation(RAW_DIR / song["file"])

        texts: list[str] = []
        for i, slide in enumerate(prs.slides):
            if i in song.get("drop", set()):
                continue
            shape_texts = [
                shape.text_frame.text for shape in slide.shapes if shape.has_text_frame
            ]
            text = "\x0b".join(shape_texts)
            if i in song.get("strip_last_line", set()):
                parts = lines_of(text)
                text = "\x0b".join(parts[:-1]) if len(parts) > 1 else text
            texts.append(text)

        unique_index_by_key: dict[str, int] = {}
        unique_lines: list[list[str]] = []
        order: list[str] = ["title"]
        text_to_key: dict[str, str] = {}

        for text in texts:
            if text in text_to_key:
                order.append(text_to_key[text])
            else:
                key = f"slide{len(unique_index_by_key) + 1}"
                unique_index_by_key[key] = len(unique_lines) + 1  # +1: title이 index 0
                unique_lines.append(lines_of(text))
                text_to_key[text] = key
                order.append(key)

        slide_xmls = [title_slide(title)] + [
            lyric_slide_from_lines(lines) for lines in unique_lines
        ]

        out_dir = SONGS_DIR / song["slug"]
        out_dir.mkdir(exist_ok=True)

        pptx_bytes = build_song_pptx(TEMPLATE, slide_xmls)
        (out_dir / "song.pptx").write_bytes(pptx_bytes)

        song_json = {
            "title": title,
            "sourceFile": "song.pptx",
            "slides": [{"key": "title", "index": 0, "label": "title"}]
            + [
                {"key": key, "index": index, "label": key}
                for key, index in unique_index_by_key.items()
            ],
            "order": order,
        }
        (out_dir / "song.json").write_text(
            json.dumps(song_json, ensure_ascii=False, indent=2), encoding="utf-8"
        )

        print(
            f"{song['slug']} ({title}): {len(slide_xmls)} unique slides "
            f"(incl. title), {len(order)} in order"
        )


if __name__ == "__main__":
    main()
